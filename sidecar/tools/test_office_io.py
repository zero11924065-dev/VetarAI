# VetarAI - Local-first multi-agent orchestration application
# Copyright (C) 2026 zero11924065-dev
# GPL-3.0-or-later（见仓库 LICENSE）
"""0.4.6：Office 输入/输出集成测试。

覆盖：
  D1 四种输入解析（docx/xlsx/pptx/pdf）真实构造文件后解析
  D2 pptx 解析含标题/正文/备注
  D3 create_document 四种类型端到端生成 + 反例
  D4 docx/xlsx/pptx 生成的文件可被对应解析器读回（闭环验证）

运行：.venv/bin/python -m sidecar.tools.test_office_io
"""
import asyncio
import io
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

PASS, FAIL = 0, 0
FAILURES = []


def check(name, cond, detail=""):
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"PASS  {name}")
    else:
        FAIL += 1
        FAILURES.append(name)
        print(f"FAIL  {name}  {detail}")


def _mk_docx(text: str = "测试段落内容") -> bytes:
    import docx
    d = docx.Document()
    d.add_paragraph(text)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _mk_pptx(title: str = "测试标题", body: str = "测试正文") -> bytes:
    from pptx import Presentation
    p = Presentation()
    sl = p.slides.add_slide(p.slide_layouts[1])
    sl.shapes.title.text = title
    sl.placeholders[1].text = body
    buf = io.BytesIO()
    p.save(buf)
    return buf.getvalue()


def _mk_xlsx(cell: str = "数据单元格") -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    wb.active["A1"] = cell
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _mk_pdf() -> bytes:
    from pypdf import PdfWriter
    w = PdfWriter()
    w.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    w.write(buf)
    return buf.getvalue()


def main():
    from sidecar.attachments.parser import parse_attachment
    from sidecar.tools.registry import execute

    # D1 四种输入解析
    t, k = parse_attachment("a.docx", _mk_docx())
    check("D1a docx 解析", k == "docx" and t and "测试段落内容" in t, f"{k}|{t}")
    t, k = parse_attachment("c.xlsx", _mk_xlsx())
    check("D1b xlsx 解析", k == "xlsx" and t and "数据单元格" in t, f"{k}|{t}")
    t, k = parse_attachment("d.pdf", _mk_pdf())
    check("D1c pdf 解析（空白页无文本返回None不报错）", k == "pdf" and t is None, f"{k}|{t}")

    # D2 pptx 解析（含标题/正文）
    t, k = parse_attachment("b.pptx", _mk_pptx())
    check("D2a pptx 解析kind", k == "pptx", str(k))
    check("D2b pptx 含标题", t is not None and "测试标题" in t, str(t))
    check("D2c pptx 含正文", t is not None and "测试正文" in t, str(t))

    tmp = tempfile.mkdtemp(prefix="office_io_")

    async def run():
        # D3 create_document 四种类型 + 反例
        r = await execute("create_document", {
            "path": f"{tmp}/报告.docx",
            "content": {"title": "T", "blocks": [
                {"type": "paragraph", "text": "正文"},
                {"type": "table", "rows": [["a", "b"], ["1", "2"]]}]}}, tmp)
        check("D3a docx 生成", r.get("ok") is True, str(r))

        r = await execute("create_document", {
            "path": f"{tmp}/数据.xlsx",
            "content": {"sheets": [{"name": "S", "rows": [["月份", "金额"], ["1月", "100"]]}]}}, tmp)
        check("D3b xlsx 生成", r.get("ok") is True, str(r))

        r = await execute("create_document", {
            "path": f"{tmp}/幻灯片.pptx",
            "content": {"slides": [{"title": "封面"}, {"title": "页2", "bullets": ["要点"]}]}}, tmp)
        check("D3c pptx 生成", r.get("ok") is True, str(r))

        r = await execute("create_document", {
            "path": f"{tmp}/总结.md",
            "content": {"title": "总结", "blocks": [{"type": "paragraph", "text": "正文"}]}}, tmp)
        check("D3d md 生成", r.get("ok") is True, str(r))

        r = await execute("create_document", {"path": f"{tmp}/x.docx", "content": "纯文本"}, tmp)
        check("D3e content非对象报错", r.get("ok") is False, str(r))

        r = await execute("create_document", {"path": f"{tmp}/y.xyz", "content": {"title": "t"}}, tmp)
        check("D3f 未知类型报错", r.get("ok") is False, str(r))

        # D4 闭环：生成的文件可被解析器读回
        docx_bytes = Path(f"{tmp}/报告.docx").read_bytes()
        t, k = parse_attachment("报告.docx", docx_bytes)
        check("D4a 生成的docx可解析", k == "docx" and t and "正文" in t, f"{k}|{t}")

        xlsx_bytes = Path(f"{tmp}/数据.xlsx").read_bytes()
        t, k = parse_attachment("数据.xlsx", xlsx_bytes)
        check("D4b 生成的xlsx可解析", k == "xlsx" and t and "100" in t, f"{k}|{t}")

        pptx_bytes = Path(f"{tmp}/幻灯片.pptx").read_bytes()
        t, k = parse_attachment("幻灯片.pptx", pptx_bytes)
        check("D4c 生成的pptx可解析", k == "pptx" and t and "要点" in t, f"{k}|{t}")

        # D5 page_break 分页块（案件证据文档每份证据独立起页的核心能力）
        r = await execute("create_document", {
            "path": f"{tmp}/证据.docx",
            "content": {"title": "证据", "blocks": [
                {"type": "heading", "level": 2, "text": "证据一 借据"},
                {"type": "paragraph", "text": "借据内容"},
                {"type": "page_break"},
                {"type": "heading", "level": 2, "text": "证据二 转账记录"},
                {"type": "paragraph", "text": "转账内容"},
            ]}}, tmp)
        check("D5a docx 分页块生成成功", r.get("ok") is True, str(r))
        import docx as _docx
        _d = _docx.Document(r["path"])
        _breaks = sum(1 for p in _d.paragraphs for run in p.runs
                      if "page" in run._element.xml and "w:br" in run._element.xml)
        check("D5b docx 含分页符", _breaks >= 1, f"breaks={_breaks}")

        r = await execute("create_document", {
            "path": f"{tmp}/证据.md",
            "content": {"blocks": [
                {"type": "paragraph", "text": "证据一"},
                {"type": "page_break"},
                {"type": "paragraph", "text": "证据二"}]}}, tmp)
        _md = Path(r["path"]).read_text(encoding="utf-8")
        check("D5c md 分页标记存在", "---" in _md, _md[:60])

    asyncio.run(run())

    import shutil
    shutil.rmtree(tmp, ignore_errors=True)

    print(f"\n===== 结果：{PASS} PASS / {FAIL} FAIL =====")
    if FAILURES:
        print("失败项：", "、".join(FAILURES))
        sys.exit(1)


if __name__ == "__main__":
    main()
