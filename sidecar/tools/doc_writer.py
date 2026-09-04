# VetarAI - Local-first multi-agent orchestration application
# Copyright (C) 2026 zero11924065-dev
#
# This file is part of VetarAI.
#
# VetarAI is free software: you can redistribute it and/or modify
# it under the terms of the GNU General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# VetarAI is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE. See the
# GNU General Public License for more details.
#
# You should have received a copy of the GNU General Public License
# along with VetarAI. If not, see <https://www.gnu.org/licenses/>.
"""0.4.6：Office 文档生成器（内置，非插件）。

Agent 通过 create_document 工具调用本模块，按结构化内容契约生成真实的
Word(.docx) / Excel(.xlsx) / PowerPoint(.pptx) / Markdown(.md) 文件，
保存到项目工作目录。全部本地生成，不联网。

内容契约（模型填写，JSON 对象）：
- docx / md：{"title": str, "blocks": [Block, ...]}
- xlsx：     {"sheets": [{"name": str, "rows": [[cell, ...], ...]}, ...]}
- pptx：     {"slides": [{"title": str, "bullets": [str,...], "notes": str}, ...]}

Block 类型：
- {"type": "heading",   "level": 1~4, "text": str}
- {"type": "paragraph", "text": str}
- {"type": "bullets",   "items": [str, ...]}
- {"type": "table",     "rows": [[cell, ...], ...]}（首行为表头）
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

SUPPORTED_DOC_TYPES = {"docx", "xlsx", "pptx", "md", "markdown"}


def _clean_text(v: Any) -> str:
    return str(v) if v is not None else ""


def _extract_blocks(content: dict) -> tuple[str, list[dict]]:
    title = _clean_text(content.get("title") or "")
    blocks = content.get("blocks") or []
    if not isinstance(blocks, list):
        blocks = []
    return title, [b for b in blocks if isinstance(b, dict)]


def _write_md(target: Path, content: dict) -> int:
    title, blocks = _extract_blocks(content)
    lines: list[str] = []
    if title:
        lines.append(f"# {title}\n")
    for b in blocks:
        t = b.get("type", "paragraph")
        if t == "page_break":
            # 0.4.6+：分页块（Markdown 用水平线标记，转 docx 时对应分页符）
            lines.append("\n---\n")
        elif t == "heading":
            level = max(1, min(int(b.get("level") or 2), 6))
            lines.append(f"{'#' * level} {_clean_text(b.get('text'))}\n")
        elif t == "bullets":
            for it in (b.get("items") or []):
                lines.append(f"- {_clean_text(it)}")
            lines.append("")
        elif t == "table":
            rows = b.get("rows") or []
            if rows:
                header = [_clean_text(c) for c in rows[0]]
                lines.append("| " + " | ".join(header) + " |")
                lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                for r in rows[1:]:
                    lines.append("| " + " | ".join(_clean_text(c) for c in r) + " |")
                lines.append("")
        else:
            lines.append(_clean_text(b.get("text")) + "\n")
    text = "\n".join(lines).strip() + "\n"
    target.write_text(text, encoding="utf-8")
    return len(text.encode("utf-8"))


def _set_cjk_font(d, font: str = "宋体", ascii_font: str = "Times New Roman") -> None:
    """0.4.6+：把文档默认（Normal 样式）中文设为宋体、西文设为 Times New Roman。

    python-docx 默认模板中文会回退为西文字体，在 WPS/Office 中显示不规范。
    法律文书用宋体最正式。通过 Normal 样式 + 默认 rPr 的 eastAsia 属性设置，
    使全文（含表格、标题）中文统一宋体。不依赖本机安装 Office，纯文件层操作。
    """
    from docx.oxml.ns import qn
    style = d.styles["Normal"]
    rpr = style.element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        from docx.oxml import OxmlElement
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    rfonts.set(qn("w:ascii"), ascii_font)
    rfonts.set(qn("w:hAnsi"), ascii_font)
    rfonts.set(qn("w:eastAsia"), font)
    # 默认字号：小四（12pt），法律文书常用
    sz = rpr.find(qn("w:sz"))
    if sz is None:
        from docx.oxml import OxmlElement
        sz = OxmlElement("w:sz")
        rpr.append(sz)
    sz.set(qn("w:val"), "24")  # 12pt = 24 半磅


def _setup_a4_page(d) -> None:
    """0.4.6+：A4 竖版 + 法律文书页边距（上下 2.54cm，左右 3.18cm 标准）。

    对文档中所有 section 统一设置。证据编排按 A4 竖版输出（用户规范）。
    """
    from docx.shared import Cm
    for sec in d.sections:
        sec.page_width = Cm(21.0)
        sec.page_height = Cm(29.7)
        sec.top_margin = Cm(2.54)
        sec.bottom_margin = Cm(2.54)
        sec.left_margin = Cm(3.18)
        sec.right_margin = Cm(3.18)


def _add_page_number_footer(d, fmt: str = "第{p}页 共{t}页") -> None:
    """0.4.6+：页脚整体页码（全档案连续页码）。

    用 Word 域代码 PAGE / NUMPAGES 自动生成「第 X 页 共 Y 页」，
    WPS/Word 打开即实时计算总页数，无需手工维护。对全部 section 生效。
    """
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    def _make_field_run(instr: str):
        """构造一个含域字符的运行（begin/instrText/end 三个 run）。"""
        r_begin = OxmlElement("w:r")
        fld_begin = OxmlElement("w:fldChar")
        fld_begin.set(qn("w:fldCharType"), "begin")
        r_begin.append(fld_begin)

        r_instr = OxmlElement("w:r")
        instr_text = OxmlElement("w:instrText")
        instr_text.set(qn("xml:space"), "preserve")
        instr_text.text = instr
        r_instr.append(instr_text)

        r_end = OxmlElement("w:r")
        fld_end = OxmlElement("w:fldChar")
        fld_end.set(qn("w:fldCharType"), "end")
        r_end.append(fld_end)
        return r_begin, r_instr, r_end

    def _make_text_run(text: str):
        r = OxmlElement("w:r")
        rpr = OxmlElement("w:rPr")
        sz = OxmlElement("w:sz"); sz.set(qn("w:val"), "18")  # 9pt 页脚小字
        rpr.append(sz)
        r.append(rpr)
        t = OxmlElement("w:t")
        t.text = text
        r.append(t)
        return r

    # 拆分模板：{p}=当前页 {t}=总页数，其余为字面文本
    for sec in d.sections:
        footer = sec.footer
        p = footer.paragraphs[0]
        p.alignment = 1  # WD_ALIGN_PARAGRAPH.CENTER
        # 清空已有内容
        for r in list(p.runs):
            r._element.getparent().remove(r._element)
        # 逐段构建：按 {p} / {t} 切分
        import re as _re
        pos = 0
        for m in _re.finditer(r"\{(p|t)\}", fmt):
            literal = fmt[pos:m.start()]
            if literal:
                p._element.append(_make_text_run(literal))
            if m.group(1) == "p":
                for run in _make_field_run("PAGE"):
                    p._element.append(run)
            else:
                for run in _make_field_run("NUMPAGES"):
                    p._element.append(run)
            pos = m.end()
        tail = fmt[pos:]
        if tail:
            p._element.append(_make_text_run(tail))


def _write_docx(target: Path, content: dict) -> int:
    import docx
    from docx.shared import Cm, Pt
    title, blocks = _extract_blocks(content)
    d = docx.Document()
    _set_cjk_font(d)
    _setup_a4_page(d)
    # 0.4.6+：整体页码页脚（默认开启，证据/目录类文档需要全档案连续页码）
    if content.get("page_number", True):
        fmt = content.get("page_number_format", "第{p}页 共{t}页")
        _add_page_number_footer(d, fmt)
    if title:
        d.add_heading(title, level=0)
    for b in blocks:
        t = b.get("type", "paragraph")
        if t == "page_break":
            # 0.4.6+：分页块（证据文档每份证据独立起页，目录页码可确定性统计）
            d.add_page_break()
        elif t == "heading":
            level = max(1, min(int(b.get("level") or 2), 4))
            d.add_heading(_clean_text(b.get("text")), level=level)
        elif t == "bullets":
            for it in (b.get("items") or []):
                d.add_paragraph(_clean_text(it), style="List Bullet")
        elif t == "table":
            rows = b.get("rows") or []
            if rows:
                ncols = max(len(r) for r in rows)
                tbl = d.add_table(rows=len(rows), cols=ncols)
                tbl.style = "Light Grid Accent 1"
                for i, r in enumerate(rows):
                    for j in range(ncols):
                        tbl.rows[i].cells[j].text = _clean_text(r[j]) if j < len(r) else ""
        elif t == "image":
            # 0.4.6+：图片证据块。
            # layout="single"（默认）：原文调入居中（营业执照/身份证/合同等正式文件）
            # layout="grid"：一行 3 列网格（聊天截图等多张并排）
            paths = b.get("paths") or ([b["path"]] if b.get("path") else [])
            layout = b.get("layout", "single")
            caption = b.get("caption") or ""
            width_cm = float(b.get("width_cm", 13.0))
            paths = [p for p in paths if isinstance(p, str) and Path(p).is_file()]
            if paths:
                if layout == "grid":
                    # 一行 3 列：用 3 列表格承载，每格一张图等宽
                    col_w_cm = 4.4  # (21-6.36)/3≈4.4cm，3列适配A4左右边距
                    for i in range(0, len(paths), 3):
                        chunk = paths[i:i + 3]
                        tbl = d.add_table(rows=1, cols=len(chunk))
                        for j, p in enumerate(chunk):
                            cell = tbl.rows[0].cells[j]
                            cell.text = ""
                            para = cell.paragraphs[0]
                            para.alignment = 1
                            run = para.add_run()
                            try:
                                run.add_picture(p, width=Cm(col_w_cm))
                            except Exception:
                                run.text = f"[图片加载失败: {Path(p).name}]"
                        # 表格无边框视觉：去掉表格样式
                        tbl.style = "Table Grid"
                    if caption:
                        cp = d.add_paragraph(caption)
                        cp.alignment = 1
                else:
                    for p in paths:
                        para = d.add_paragraph()
                        para.alignment = 1
                        run = para.add_run()
                        try:
                            run.add_picture(p, width=Cm(width_cm))
                        except Exception:
                            run.text = f"[图片加载失败: {Path(p).name}]"
                    if caption:
                        cp = d.add_paragraph(caption)
                        cp.alignment = 1
        else:
            d.add_paragraph(_clean_text(b.get("text")))
    d.save(str(target))
    return target.stat().st_size


def _write_xlsx(target: Path, content: dict) -> int:
    import openpyxl
    sheets = content.get("sheets") or []
    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # 移除默认空表
    if not sheets:
        wb.create_sheet("Sheet1")
    for idx, sh in enumerate(sheets):
        if not isinstance(sh, dict):
            continue
        name = _clean_text(sh.get("name") or f"Sheet{idx + 1}")[:31] or f"Sheet{idx + 1}"
        ws = wb.create_sheet(name)
        for r_idx, row in enumerate(sh.get("rows") or [], start=1):
            for c_idx, cell in enumerate(row, start=1):
                v = cell
                # 纯数字字符串转数字，便于 Excel 计算
                if isinstance(v, str):
                    s = v.strip()
                    if s.lstrip("-").isdigit():
                        v = int(s)
                    elif s.replace(".", "", 1).lstrip("-").isdigit() and s.count(".") == 1:
                        v = float(s)
                ws.cell(row=r_idx, column=c_idx, value=v)
    wb.save(str(target))
    return target.stat().st_size


def _write_pptx(target: Path, content: dict) -> int:
    from pptx import Presentation
    from pptx.util import Pt
    slides = content.get("slides") or []
    prs = Presentation()
    title_layout = prs.slide_layouts[0]   # 标题页
    bullet_layout = prs.slide_layouts[1]  # 标题+内容
    blank_title = prs.slide_layouts[5]    # 仅标题
    if not slides:
        slides = [{"title": _clean_text(content.get("title") or "演示文稿"), "bullets": []}]
    for idx, sl in enumerate(slides):
        if not isinstance(sl, dict):
            continue
        s_title = _clean_text(sl.get("title") or "")
        bullets = [b for b in (sl.get("bullets") or [])]
        notes = _clean_text(sl.get("notes") or "")
        layout = title_layout if idx == 0 and not bullets else (bullet_layout if bullets else blank_title)
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = s_title
        if bullets and layout in (bullet_layout,):
            tf = slide.placeholders[1].text_frame
            tf.text = _clean_text(bullets[0])
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = _clean_text(b)
                p.level = 0
        if notes and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = notes
        elif notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass
    prs.save(str(target))
    return target.stat().st_size


def write_document(doc_type: str, target: Path, content: dict) -> int:
    """按类型生成文档，返回写入字节数。不支持的类型抛 ValueError。"""
    if not isinstance(content, dict):
        raise ValueError("bad_arg: content 必须是结构化 JSON 对象（非纯文本）")
    dt = (doc_type or "").lower().strip()
    if dt in ("md", "markdown"):
        return _write_md(target, content)
    if dt == "docx":
        return _write_docx(target, content)
    if dt in ("xlsx", "xls"):
        return _write_xlsx(target, content)
    if dt in ("pptx", "ppt"):
        return _write_pptx(target, content)
    raise ValueError(f"bad_arg: doc_type '{doc_type}' 不支持（可选 docx/xlsx/pptx/md）")
